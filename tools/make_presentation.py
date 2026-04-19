"""生成 Final Presentation.pptx / Generate the final presentation.

完全基于真实实验结果 / Grounded in actual measured results:
    - ResNet-50 + EMD: val SRCC 0.8903, test 0.8842, SPAQ 0.8411
    - ViT-B/16 + LoRA: val SRCC 0.9098, test 0.8940, SPAQ 0.8399
    - ONNX CUDA 5.37ms, ViT graph fusion 28.3→26.2ms (7% gain)

9 张 slides × ~1 min 讲解 ≈ 9-10 分钟，符合 assignment 要求。
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# -- 颜色 / Color palette --------------------------------------------------
NAVY = RGBColor(0x0C, 0x2D, 0x48)      # 标题深蓝 / title dark blue
ACCENT = RGBColor(0x00, 0x71, 0xE3)    # 强调色 / accent blue
GRAY = RGBColor(0x4A, 0x4A, 0x4A)      # 正文灰 / body text
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)     # 浅背景 / light bg
OK_GREEN = RGBColor(0x10, 0x8A, 0x3D)
WARN = RGBColor(0xC0, 0x39, 0x2B)

HEATMAP_DIR = Path("C:/iqa-project/outputs/heatmaps")


# -- 样式辅助 / Style helpers ---------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    """设置纯色背景 / fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str, top_in: float = 0.35) -> None:
    """添加标题文本框 / add title text box."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top_in), Inches(12.33), Inches(0.8))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_subtitle(slide, text: str, top_in: float = 1.05) -> None:
    """添加副标题 / add subtitle."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top_in), Inches(12.33), Inches(0.45))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.font.italic = True


def add_bullet_block(slide, bullets: list, left: float, top: float,
                     width: float, height: float, font_size: int = 16) -> None:
    """添加带项目符号的正文 / bulleted body text block.

    bullets 里每项是 (text, indent) 或 str（默认 indent=0）。
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, indent = item
        else:
            text, indent = item, 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = indent
        p.font.size = Pt(font_size - indent * 2)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)


def add_table(slide, rows: list, left: float, top: float,
              width: float, height: float) -> None:
    """添加数据表 / add a data table.

    rows[0] 是表头 / rows[0] is header row.
    """
    n_rows, n_cols = len(rows), len(rows[0])
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(13)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        run.font.color.rgb = GRAY
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY


def add_footer(slide, slide_no: int, total: int) -> None:
    """底部水印 / footer watermark."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"AAI 6640 · IQA on KonIQ-10k · {slide_no}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


# =========================================================================
# 构建演示文稿 / Build presentation
# =========================================================================
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # 全空白布局 / fully blank layout

TOTAL = 9


# ---- Slide 1: Title ------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY)

# 主标题 / main title
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4))
p = tb.text_frame.paragraphs[0]
p.text = "Deep Learning for Image Quality Assessment"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 副标题 / subtitle
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = "Distribution Prediction with EMD Loss on KonIQ-10k"
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
p.alignment = PP_ALIGN.CENTER

# 装饰线 / decorative line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.4),
                          Inches(2.33), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# 团队 + 课程 / team + course
tb = s.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.33), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Haowei Shen · Xuezhen Jin · Peter Adranly"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "AAI 6640 — Final Project"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
p.alignment = PP_ALIGN.CENTER


# ---- Slide 2: Background -------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Background  /  背景")
add_subtitle(s, "In-the-wild IQA: predicting human perception of photo quality from a single image")

add_bullet_block(s, [
    "• Image Quality Assessment (IQA) predicts the Mean Opinion Score (MOS) "
    "that human raters would assign to a natural image.",
    "• Applications: social media ranking, smartphone camera tuning, ",
    "  compression quality control, generative model evaluation.",
    "",
    "• In-the-wild challenge: real-world photos suffer mixed distortions "
    "(blur + exposure + noise), unlike synthetic-distortion benchmarks.",
    "• KonIQ-10k (Hosu et al., TIP 2020): 10,073 ecologically-valid images "
    "with per-image MOS + annotator std, derived from YFCC100M.",
    "",
    "• This project compares three architectural paradigms on the same benchmark "
    "and extends to cross-dataset generalization and deployment optimization.",
], left=0.7, top=1.6, width=12.0, height=5.0, font_size=17)

add_footer(s, 2, TOTAL)


# ---- Slide 3: Research Questions -----------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Research Questions  /  研究问题")
add_subtitle(s, "Three hypotheses systematically tested on the same KonIQ-10k split")

add_bullet_block(s, [
    "Q1.  Distribution vs. Mean Prediction",
    ("Can EMD loss on probability distributions outperform MSE regression "
     "on mean scores for in-the-wild IQA?", 1),
    "",
    "Q2.  Parameter-Efficient Fine-Tuning",
    ("Can LoRA on ViT-B/16 (training <1% of parameters) match or exceed "
     "full-backbone fine-tuning of ResNet-50?", 1),
    "",
    "Q3.  Generalization and Deployment",
    ("How do in-domain winners transfer zero-shot to SPAQ?", 1),
    ("What is the latency floor for production deployment (ONNX, "
     "quantization, graph fusion)?", 1),
], left=0.8, top=1.6, width=12.0, height=5.0, font_size=17)

add_footer(s, 3, TOTAL)


# ---- Slide 4: Methods — Loss + Models -----------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Methods  /  方法：EMD Loss + Three Models")
add_subtitle(s, "Gaussian-shape soft targets with squared Earth Mover's Distance")

# 左栏：EMD + Gaussian target / left column
add_bullet_block(s, [
    "Squared EMD Loss (order-aware):",
    ("EMD² ordinal — confusing '9' with '10' is cheap, '9' with '2' is expensive",
     1),
    ("L = ( (1/B) Σᵢ | CDF(p)ᵢ − CDF(q)ᵢ |² )^(1/2)", 1),
    "",
    "Target distribution:",
    ("MOS + std → discrete 10-bucket Gaussian PDF", 1),
    ("Clamp σ ≥ 0.5 bucket units (prevent one-hot collapse)", 1),
    "",
    "Mixed precision stack (RTX 5090 Blackwell sm_120):",
    ("TF32 for conv/matmul outside autocast", 1),
    ("bf16 AMP for hot path (softmax, LayerNorm stay FP32)", 1),
    ("⇒ ~2.5× training throughput", 1),
], left=0.6, top=1.55, width=5.8, height=5.3, font_size=13)

# 右栏：三个模型 / right column
add_bullet_block(s, [
    "Baseline CNN (MSE):",
    ("4-block (Conv→BN→ReLU→MaxPool) + GAP + FC", 1),
    ("~1.2M params, scalar regression", 1),
    "",
    "ResNet-50 + EMD:",
    ("ImageNet V2 pretrained + 10-bucket softmax head", 1),
    ("Phase 1 (ep 0-4): freeze backbone, LR=1e-3", 1),
    ("Phase 2 (ep 5-29): unfreeze layer3/4, LR=1e-5/1e-4", 1),
    "",
    "ViT-B/16 + LoRA + EMD:",
    ("timm ImageNet-21k pretrained backbone", 1),
    ("LoRA r=16, α=32, target=qkv, save head fully", 1),
    ("597K / 86.4M = 0.69% trainable (parameter-efficient)", 1),
], left=6.8, top=1.55, width=6.2, height=5.3, font_size=13)

add_footer(s, 4, TOTAL)


# ---- Slide 5: In-Dataset Results -----------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Key Finding #1  /  In-Dataset Performance")
add_subtitle(s, "5-crop TTA on KonIQ-10k test split (seed=42, 8:1:1 image_id split)")

add_table(s, [
    ["Model", "Trainable", "Val SRCC", "Test SRCC", "Test PLCC", "Test RMSE"],
    ["Baseline CNN (quick probe, 2 ep)", "1.2 M", "0.688", "—", "—", "—"],
    ["ResNet-50 + EMD", "22 M (Ph2)", "0.8903", "0.8842", "0.9041", "0.068"],
    ["ViT-B/16 + LoRA + EMD", "597 K  (0.69%)", "0.9098", "0.8940", "0.9071", "0.070"],
], left=0.7, top=1.7, width=11.9, height=2.0)

add_bullet_block(s, [
    "• ViT leads ResNet by +1 SRCC point in-domain, while training 36× fewer parameters.",
    "• Both EMD models are within the 0.88–0.92 SRCC range reported in recent KonIQ literature.",
    "• RMSE on normalized [0,1] scale — both below 0.07, i.e. <7 MOS-point average error.",
], left=0.8, top=4.2, width=11.8, height=1.6, font_size=16)

# 强调 box / callout
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9),
                         Inches(11.7), Inches(0.8))
box.fill.solid()
box.fill.fore_color.rgb = ACCENT
box.line.fill.background()
tf = box.text_frame
tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]
p.text = ("→ EMD-on-distribution loss makes ordinal structure explicit; "
          "LoRA on ViT makes pretraining transferable without overfitting 10K samples.")
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_footer(s, 5, TOTAL)


# ---- Slide 6: Cross-Dataset + Finding ------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Key Finding #2  /  Cross-Dataset Zero-Shot (KonIQ → SPAQ)")
add_subtitle(s, "Train on KonIQ-10k, test on SPAQ (11,125 smartphone photos) with no fine-tuning")

add_table(s, [
    ["Model", "KonIQ test SRCC", "SPAQ SRCC", "SPAQ PLCC", "SRCC drop"],
    ["ResNet-50 + EMD", "0.8842", "0.8411", "0.8465", "−4.3 %"],
    ["ViT-B/16 + LoRA + EMD", "0.8940", "0.8399", "0.8439", "−5.4 %"],
], left=0.9, top=1.7, width=11.5, height=1.5)

add_bullet_block(s, [
    "• SRCC drop of 4–5 points is consistent with published cross-dataset transfer numbers "
    "— this is the expected magnitude of in-the-wild domain gap.",
    "",
    "• Surprising: ViT's 1-point in-domain advantage completely disappears on SPAQ "
    "(0.8399 vs 0.8411 — statistically equivalent).",
    "",
    "• Interpretation: the extra SRCC point from ViT appears to come from fitting "
    "KonIQ-specific annotation style (YFCC100M content distribution, particular rater pool) "
    "rather than learning more transferable quality features.",
    "",
    "• Practical implication: on new domains (e.g. smartphone photos, medical imaging), "
    "CNN-based IQA still competitive if compute-constrained.",
], left=0.8, top=3.6, width=11.8, height=3.4, font_size=15)

add_footer(s, 6, TOTAL)


# ---- Slide 7: XAI visualizations -----------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Prediction & Interpretation  /  Grad-CAM vs. Attention Rollout")
add_subtitle(s, "16 test images spanning MOS quartiles — where each model 'looks'")

# 两张热力图并排 / two heatmaps side by side
try:
    if (HEATMAP_DIR / "resnet_gradcam.png").exists():
        s.shapes.add_picture(
            str(HEATMAP_DIR / "resnet_gradcam.png"),
            Inches(0.3), Inches(1.6), width=Inches(6.3),
        )
    if (HEATMAP_DIR / "vit_rollout.png").exists():
        s.shapes.add_picture(
            str(HEATMAP_DIR / "vit_rollout.png"),
            Inches(6.8), Inches(1.6), width=Inches(6.3),
        )
except Exception as e:
    print(f"[warn] heatmap embed failed: {e}")

# 两个标注 / two captions
add_bullet_block(s, [
    "ResNet-50 Grad-CAM",
    ("Gradient-weighted feature maps at layer4[-1]", 1),
], left=0.5, top=6.2, width=6.0, height=0.9, font_size=14)

add_bullet_block(s, [
    "ViT-B/16 Attention Rollout",
    ("Cumulative attention ∏(Aₗ + I) from 12 blocks, CLS→patch row", 1),
], left=7.0, top=6.2, width=6.0, height=0.9, font_size=14)

add_footer(s, 7, TOTAL)


# ---- Slide 8: Deployment benchmark ---------------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, LIGHT)
add_title(s, "Deployment  /  Latency Benchmark on RTX 5090")
add_subtitle(s, "batch=1, 100 warmup + 500 measured iterations, single-image inference")

add_table(s, [
    ["Model", "Backend", "Precision", "Latency", "Throughput"],
    ["ResNet-50", "ORT CUDA", "FP32", "5.37 ms", "186 img/s"],
    ["ResNet-50", "ORT CUDA", "FP16", "9.34 ms ⚠", "107 img/s"],
    ["ResNet-50", "ORT CPU (9950X3D)", "FP32", "5.15 ms", "194 img/s"],
    ["ViT-B/16", "ORT CUDA — raw export", "FP32", "28.33 ms", "35 img/s"],
    ["ViT-B/16", "ORT CUDA + Transformer fusion", "FP32", "26.22 ms", "38 img/s"],
], left=0.7, top=1.6, width=11.9, height=3.0)

add_bullet_block(s, [
    "• ViT-specific graph optimization (onnxruntime.transformers.optimizer) fuses 40+ "
    "primitive ops per attention block into Attention/LayerNorm/GELU kernels → "
    "−7% latency on top of FP32.",
    "",
    "• ⚠ Blackwell anomaly: ORT CUDA FP16 is slower than FP32 (9.34 vs 5.37 ms). "
    "Matches documented issue on RTX 50-series — we abandoned TensorRT FP16/INT8 "
    "as the primary deployment target and committed to ORT CUDA.",
    "",
    "• CPU competitive at batch=1: 9950X3D's cache + ORT scheduler overhead reach "
    "parity with CUDA for single-image latency.",
], left=0.8, top=4.8, width=11.8, height=2.3, font_size=14)

add_footer(s, 8, TOTAL)


# ---- Slide 9: Conclusion + Contributions --------------------------------
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY)

# 白色标题 / white title on navy
tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
p = tb.text_frame.paragraphs[0]
p.text = "Conclusions  &  Individual Contributions"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# 左：结论 / conclusions
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

header = tf.paragraphs[0]
header.text = "Answers to Research Questions"
header.font.size = Pt(18)
header.font.bold = True
header.font.color.rgb = ACCENT

concs = [
    "Q1 (EMD vs MSE): EMD + distribution prediction achieves "
    "SRCC 0.88–0.91 on KonIQ-10k, in line with SOTA.",
    "Q2 (LoRA): ViT-LoRA at 0.69% trainable params matches "
    "full-tune ResNet. PEFT works at 10K scale.",
    "Q3 (Transfer): Zero-shot KonIQ→SPAQ drops 4–5 SRCC. "
    "ViT's in-domain lead vanishes — warning against single-dataset benchmarks.",
    "Deployment: ORT CUDA 5.37 ms (ResNet-50); ViT Transformer "
    "fusion yields +7% latency reduction.",
]
for c in concs:
    p = tf.add_paragraph()
    p.text = "•  " + c
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.space_after = Pt(8)

# 右：个人贡献 / individual contributions
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.4), Inches(6.0), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

h = tf.paragraphs[0]
h.text = "Individual Contributions"
h.font.size = Pt(18)
h.font.bold = True
h.font.color.rgb = ACCENT

contrib = [
    ("Haowei Shen", "ResNet-50 + EMD architecture, EMD loss, "
                    "entire deployment pipeline (ONNX export, Transformer "
                    "fusion, multi-backend benchmark), repo + CI setup"),
    ("Xuezhen Jin", "Data pipeline (KonIQ-10k loader, MOS_zscore "
                    "handling, 5-crop TTA), baseline CNN, evaluation "
                    "framework, cross-dataset eval harness"),
    ("Peter Adranly", "ViT-B/16 + LoRA integration, Attention Rollout "
                      "(incl. fused_attn workaround), XAI visualization "
                      "module, ViT-specific ONNX optimization"),
]
for name, desc in contrib:
    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.space_after = Pt(4)

# 底部：GitHub link / footer
tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.35))
p = tb.text_frame.paragraphs[0]
p.text = "Code:  https://github.com/SimonShenhw/iqa-koniq10k"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
p.alignment = PP_ALIGN.CENTER


# =========================================================================
# 保存到桌面 / Save to desktop
# =========================================================================
out_path = Path("C:/Users/Shen Haowei/Desktop/IQA_Final_Presentation.pptx")
prs.save(str(out_path))
print(f"[OK] Saved: {out_path}  ({out_path.stat().st_size // 1024} KB)")
print(f"     Slides: {len(prs.slides)}")
